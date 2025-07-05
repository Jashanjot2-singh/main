import streamlit as st

st.set_page_config(page_title="🤖 Agentic RAG Chatbot with MCP", layout="wide")
from datetime import datetime
from PIL import Image
import pytesseract
from langchain.schema import Document
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.prompts import ChatPromptTemplate
from qdrant_client import QdrantClient
from qdrant_client.http.models import PointStruct, VectorParams, Distance
from docx import Document as DocxDocument
import tempfile
import os
import pandas as pd
import requests
import json
import uuid
import time
from typing import Dict, List, Any, Optional
from dataclasses import dataclass, asdict
from pptx import Presentation
import threading
from queue import Queue
import traceback


# ==================== MCP (Model Context Protocol) Framework ====================

@dataclass
class MCPMessage:
    """Model Context Protocol Message Structure"""
    sender: str
    receiver: str
    type: str
    trace_id: str
    payload: Dict[str, Any]
    timestamp: float = None

    def __post_init__(self):
        if self.timestamp is None:
            self.timestamp = time.time()

    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)


class MCPBus:
    """In-memory Message Bus for MCP communication"""

    def __init__(self):
        self.message_queue = Queue()
        self.subscribers = {}
        self.message_history = []
        self.running = True

    def publish(self, message: MCPMessage):
        """Publish message to the bus"""
        self.message_queue.put(message)
        self.message_history.append(message)

    def subscribe(self, agent_name: str, callback):
        """Subscribe agent to receive messages"""
        self.subscribers[agent_name] = callback

    def process_messages(self):
        """Process messages in the queue"""
        while self.running:
            try:
                if not self.message_queue.empty():
                    message = self.message_queue.get(timeout=1)
                    if message.receiver in self.subscribers:
                        self.subscribers[message.receiver](message)
                time.sleep(0.1)
            except:
                continue


# Global MCP Bus instance
mcp_bus = MCPBus()


# ==================== Base Agent Class ====================

class BaseAgent:
    """Base class for all agents"""

    def __init__(self, name: str):
        self.name = name
        self.mcp_bus = mcp_bus
        self.mcp_bus.subscribe(self.name, self.handle_message)

    def send_message(self, receiver: str, message_type: str, payload: Dict[str, Any], trace_id: str = None):
        """Send MCP message to another agent"""
        if trace_id is None:
            trace_id = str(uuid.uuid4())

        message = MCPMessage(
            sender=self.name,
            receiver=receiver,
            type=message_type,
            trace_id=trace_id,
            payload=payload
        )

        self.mcp_bus.publish(message)
        return trace_id

    def handle_message(self, message: MCPMessage):
        """Handle incoming MCP messages"""
        pass


# ==================== Document Processing Functions ====================

def ocr_image(image_path):
    """Extract text from images using OCR"""
    try:
        # Adjust path as needed for your system
        pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        image = Image.open(image_path)
        return pytesseract.image_to_string(image)
    except Exception as e:
        return f"OCR Error: {str(e)}"


def load_pdf(path):
    """Load PDF documents"""
    try:
        return PyPDFLoader(path).load()
    except Exception as e:
        return [Document(page_content=f"PDF Load Error: {str(e)}")]


def load_txt(path):
    """Load text documents"""
    try:
        with open(path, encoding="utf-8") as f:
            return [Document(page_content=f.read())]
    except Exception as e:
        return [Document(page_content=f"Text Load Error: {str(e)}")]


def load_docx(path):
    """Load DOCX documents"""
    try:
        doc = DocxDocument(path)
        content = "\n".join([p.text for p in doc.paragraphs])
        return [Document(page_content=content)]
    except Exception as e:
        return [Document(page_content=f"DOCX Load Error: {str(e)}")]


def load_pptx(path):
    """Load PPTX presentations"""
    try:
        prs = Presentation(path)
        content = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"Slide {slide_num}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            content.append(slide_text)
        return [Document(page_content="\n".join(content))]
    except Exception as e:
        return [Document(page_content=f"PPTX Load Error: {str(e)}")]


def load_csv(path):
    """Load CSV files"""
    try:
        df = pd.read_csv(path)
        return df
    except Exception as e:
        return pd.DataFrame({"error": [f"CSV Load Error: {str(e)}"]})


def to_docs_from_csv(df):
    """Convert CSV DataFrame to Documents"""
    try:
        # Convert each row to a document
        docs = []
        headers = df.columns.tolist()

        # Add headers as context
        docs.append(Document(page_content=f"CSV Headers: {', '.join(headers)}"))

        # Add each row as a document
        for idx, row in df.iterrows():
            row_content = ", ".join([f"{col}: {val}" for col, val in row.items()])
            docs.append(Document(page_content=f"Row {idx + 1}: {row_content}"))

        return docs
    except Exception as e:
        return [Document(page_content=f"CSV Conversion Error: {str(e)}")]


# ==================== Agent Implementations ====================

class IngestionAgent(BaseAgent):
    """Agent responsible for document ingestion and preprocessing"""

    def __init__(self):
        super().__init__("IngestionAgent")

    def handle_message(self, message: MCPMessage):
        """Handle incoming messages"""
        if message.type == "INGESTION_REQUEST":
            self.process_documents(message)

    def process_documents(self, message: MCPMessage):
        """Process uploaded documents"""
        try:
            files_data = message.payload.get("files", [])
            processed_docs = []

            for file_data in files_data:
                file_path = file_data["path"]
                file_ext = file_data["ext"]
                file_name = file_data["name"]

                # Process based on file type
                if file_ext in ["png", "jpg", "jpeg"]:
                    content = ocr_image(file_path)
                    docs = [Document(page_content=content, metadata={"source": file_name})]
                elif file_ext == "pdf":
                    docs = load_pdf(file_path)
                elif file_ext == "txt":
                    docs = load_txt(file_path)
                elif file_ext == "docx":
                    docs = load_docx(file_path)
                elif file_ext == "pptx":
                    docs = load_pptx(file_path)
                elif file_ext == "csv":
                    df = load_csv(file_path)
                    docs = to_docs_from_csv(df)
                else:
                    docs = [Document(page_content=f"Unsupported file type: {file_ext}")]

                # Add metadata
                for doc in docs:
                    if not hasattr(doc, 'metadata') or doc.metadata is None:
                        doc.metadata = {}
                    doc.metadata.update({
                        "source": file_name,
                        "file_type": file_ext,
                        "processed_at": datetime.now().isoformat()
                    })

                processed_docs.extend(docs)

            # Send processed documents to RetrievalAgent
            self.send_message(
                receiver="RetrievalAgent",
                message_type="INGESTION_COMPLETE",
                payload={
                    "documents": [{"content": doc.page_content, "metadata": doc.metadata} for doc in processed_docs],
                    "status": "success"
                },
                trace_id=message.trace_id
            )

        except Exception as e:
            self.send_message(
                receiver="CoordinatorAgent",
                message_type="INGESTION_ERROR",
                payload={"error": str(e), "traceback": traceback.format_exc()},
                trace_id=message.trace_id
            )


class RetrievalAgent(BaseAgent):
    """Agent responsible for embedding generation and retrieval"""

    def __init__(self):
        super().__init__("RetrievalAgent")
        self.vector_store = {}
        self.embeddings_cache = {}

    def handle_message(self, message: MCPMessage):
        """Handle incoming messages"""
        if message.type == "INGESTION_COMPLETE":
            self.index_documents(message)
        elif message.type == "RETRIEVAL_REQUEST":
            self.retrieve_context(message)

    def compute_embeddings(self, texts: List[str]) -> List[List[float]]:
        """Compute embeddings using Ollama"""
        try:
            embeddings = []
            for text in texts:
                # Check cache first
                if text in self.embeddings_cache:
                    embeddings.append(self.embeddings_cache[text])
                    continue

                response = requests.post(
                    "http://localhost:11434/api/embeddings",
                    json={"model": "nomic-embed-text", "prompt": text},
                    timeout=30
                )

                if response.status_code == 200:
                    embedding = response.json().get("embedding", [])
                    self.embeddings_cache[text] = embedding
                    embeddings.append(embedding)
                else:
                    embeddings.append([0.0] * 768)  # Default embedding size

            return embeddings
        except Exception as e:
            st.error(f"Embedding error: {str(e)}")
            return [[0.0] * 768 for _ in texts]

    def index_documents(self, message: MCPMessage):
        """Index documents in vector store"""
        try:
            documents = message.payload.get("documents", [])

            if not documents:
                return

            # Split documents into chunks
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=500,
                chunk_overlap=50
            )

            chunks = []
            for doc in documents:
                doc_obj = Document(page_content=doc["content"], metadata=doc["metadata"])
                doc_chunks = text_splitter.split_documents([doc_obj])
                chunks.extend(doc_chunks)

            # Extract text content
            texts = [chunk.page_content for chunk in chunks]

            if not texts:
                return

            # Compute embeddings
            embeddings = self.compute_embeddings(texts)

            # Store in Qdrant
            collection_name = f"docs_{message.trace_id}"
            self.create_qdrant_collection(collection_name, len(embeddings[0]) if embeddings else 768)
            self.insert_into_qdrant(collection_name, texts, embeddings, chunks)

            # Store collection info
            self.vector_store[message.trace_id] = {
                "collection_name": collection_name,
                "document_count": len(texts),
                "indexed_at": datetime.now().isoformat()
            }

            # Notify completion
            self.send_message(
                receiver="CoordinatorAgent",
                message_type="INDEXING_COMPLETE",
                payload={
                    "collection_name": collection_name,
                    "document_count": len(texts),
                    "status": "success"
                },
                trace_id=message.trace_id
            )

        except Exception as e:
            self.send_message(
                receiver="CoordinatorAgent",
                message_type="INDEXING_ERROR",
                payload={"error": str(e), "traceback": traceback.format_exc()},
                trace_id=message.trace_id
            )

    def retrieve_context(self, message: MCPMessage):
        """Retrieve relevant context for a query"""
        try:
            query = message.payload.get("query", "")
            trace_id = message.trace_id

            if trace_id not in self.vector_store:
                raise ValueError("No indexed documents found for this session")

            collection_name = self.vector_store[trace_id]["collection_name"]

            # Compute query embedding
            query_embedding = self.compute_embeddings([query])[0]

            # Search in Qdrant
            relevant_chunks = self.search_qdrant(collection_name, query_embedding, limit=5)

            # Send results to LLMResponseAgent
            self.send_message(
                receiver="LLMResponseAgent",
                message_type="RETRIEVAL_RESULT",
                payload={
                    "retrieved_context": relevant_chunks,
                    "query": query,
                    "collection_name": collection_name
                },
                trace_id=trace_id
            )

        except Exception as e:
            self.send_message(
                receiver="CoordinatorAgent",
                message_type="RETRIEVAL_ERROR",
                payload={"error": str(e), "traceback": traceback.format_exc()},
                trace_id=message.trace_id
            )

    def create_qdrant_collection(self, collection_name: str, vector_size: int):
        """Create Qdrant collection"""
        try:
            client = QdrantClient(host="localhost", port=6333)
            client.recreate_collection(
                collection_name=collection_name,
                vectors_config=VectorParams(size=vector_size, distance=Distance.COSINE)
            )
        except Exception as e:
            st.error(f"Qdrant collection creation error: {str(e)}")

    def insert_into_qdrant(self, collection_name: str, texts: List[str], embeddings: List[List[float]],
                           chunks: List[Document]):
        """Insert documents into Qdrant"""
        try:
            client = QdrantClient(host="localhost", port=6333)
            points = []

            for i, (text, embedding, chunk) in enumerate(zip(texts, embeddings, chunks)):
                point = PointStruct(
                    id=i,
                    vector=embedding,
                    payload={
                        "text": text,
                        "metadata": chunk.metadata if hasattr(chunk, 'metadata') else {}
                    }
                )
                points.append(point)

            client.upsert(collection_name=collection_name, points=points)

        except Exception as e:
            st.error(f"Qdrant insertion error: {str(e)}")

    def search_qdrant(self, collection_name: str, query_vector: List[float], limit: int = 5) -> List[str]:
        """Search for relevant documents in Qdrant"""
        try:
            client = QdrantClient(host="localhost", port=6333)
            results = client.search(
                collection_name=collection_name,
                query_vector=query_vector,
                limit=limit
            )

            return [result.payload.get("text", "") for result in results]

        except Exception as e:
            st.error(f"Qdrant search error: {str(e)}")
            return []


class LLMResponseAgent(BaseAgent):
    """Agent responsible for generating LLM responses"""

    def __init__(self):
        super().__init__("LLMResponseAgent")

    def handle_message(self, message: MCPMessage):
        """Handle incoming messages"""
        if message.type == "RETRIEVAL_RESULT":
            self.generate_response(message)

    def generate_response(self, message: MCPMessage):
        """Generate response using LLM"""
        try:
            query = message.payload.get("query", "")
            context_chunks = message.payload.get("retrieved_context", [])

            # Combine context
            context = "\n\n".join(context_chunks)

            # Create prompt
            prompt_template = ChatPromptTemplate.from_template("""
            Answer the question based only on the following context:

            Context:
            {context}

            Question: {question}

            Provide a detailed and well-supported answer. If the context doesn't contain enough information to answer the question, say so clearly.
            """)

            prompt = prompt_template.format(context=context, question=query)

            # Generate response using Ollama
            response = requests.post(
                "http://localhost:11434/api/generate",
                json={
                    "model": "llama3.2:1b",
                    "prompt": prompt,
                    "stream": False
                }
            )

            if response.status_code == 200:
                llm_response = response.json().get("response", "No response generated")
            else:
                llm_response = f"Error generating response: {response.status_code}"

            # Send final response
            self.send_message(
                receiver="CoordinatorAgent",
                message_type="RESPONSE_COMPLETE",
                payload={
                    "response": llm_response,
                    "context_chunks": context_chunks,
                    "query": query,
                    "status": "success"
                },
                trace_id=message.trace_id
            )

        except Exception as e:
            self.send_message(
                receiver="CoordinatorAgent",
                message_type="RESPONSE_ERROR",
                payload={"error": str(e), "traceback": traceback.format_exc()},
                trace_id=message.trace_id
            )


class CoordinatorAgent(BaseAgent):
    """Coordinator agent that orchestrates the entire RAG pipeline"""

    def __init__(self):
        super().__init__("CoordinatorAgent")
        self.session_states = {}

    def handle_message(self, message: MCPMessage):
        """Handle incoming messages"""
        if message.type == "INGESTION_COMPLETE":
            self.handle_ingestion_complete(message)
        elif message.type == "INDEXING_COMPLETE":
            self.handle_indexing_complete(message)
        elif message.type == "RESPONSE_COMPLETE":
            self.handle_response_complete(message)
        elif message.type in ["INGESTION_ERROR", "INDEXING_ERROR", "RETRIEVAL_ERROR", "RESPONSE_ERROR"]:
            self.handle_error(message)

    def process_documents(self, files_data: List[Dict], trace_id: str = None):
        """Start document processing pipeline"""
        if trace_id is None:
            trace_id = str(uuid.uuid4())

        self.session_states[trace_id] = {
            "status": "processing",
            "stage": "ingestion",
            "started_at": datetime.now().isoformat()
        }

        # Send to IngestionAgent
        self.send_message(
            receiver="IngestionAgent",
            message_type="INGESTION_REQUEST",
            payload={"files": files_data},
            trace_id=trace_id
        )

        return trace_id

    def process_query(self, query: str, trace_id: str):
        """Process user query"""
        if trace_id not in self.session_states:
            return {"error": "No documents indexed for this session"}

        # Send to RetrievalAgent
        self.send_message(
            receiver="RetrievalAgent",
            message_type="RETRIEVAL_REQUEST",
            payload={"query": query},
            trace_id=trace_id
        )

        return {"status": "processing", "trace_id": trace_id}

    def handle_ingestion_complete(self, message: MCPMessage):
        """Handle ingestion completion"""
        trace_id = message.trace_id
        self.session_states[trace_id]["stage"] = "indexing"

    def handle_indexing_complete(self, message: MCPMessage):
        """Handle indexing completion"""
        trace_id = message.trace_id
        self.session_states[trace_id].update({
            "status": "ready",
            "stage": "ready",
            "collection_name": message.payload.get("collection_name"),
            "document_count": message.payload.get("document_count")
        })

    def handle_response_complete(self, message: MCPMessage):
        """Handle response completion"""
        trace_id = message.trace_id
        self.session_states[trace_id]["last_response"] = message.payload

    def handle_error(self, message: MCPMessage):
        """Handle errors"""
        trace_id = message.trace_id
        self.session_states[trace_id].update({
            "status": "error",
            "error": message.payload.get("error"),
            "error_type": message.type
        })

    def get_session_status(self, trace_id: str) -> Dict:
        """Get session status"""
        return self.session_states.get(trace_id, {"status": "not_found"})


# ==================== Initialize Agents ====================

@st.cache_resource
def initialize_agents():
    """Initialize all agents"""
    coordinator = CoordinatorAgent()
    ingestion_agent = IngestionAgent()
    retrieval_agent = RetrievalAgent()
    llm_response_agent = LLMResponseAgent()

    # Start MCP bus processing in background
    threading.Thread(target=mcp_bus.process_messages, daemon=True).start()

    return coordinator, ingestion_agent, retrieval_agent, llm_response_agent


# Initialize agents
coordinator, ingestion_agent, retrieval_agent, llm_response_agent = initialize_agents()

# ==================== Streamlit UI ====================

# Initialize session state
if "files" not in st.session_state:
    st.session_state.files = []
if "trace_id" not in st.session_state:
    st.session_state.trace_id = None
if "query_results" not in st.session_state:
    st.session_state.query_results = {}
if "mcp_messages" not in st.session_state:
    st.session_state.mcp_messages = []

# Title and description
st.title("🤖 Agentic RAG Chatbot with MCP")
st.markdown("### Multi-Agent Document QA System using Model Context Protocol")

# Sidebar for system status
with st.sidebar:
    st.header("🔧 System Status")

    # Check Ollama connection
    try:
        ollama_response = requests.get("http://localhost:11434/api/tags", timeout=5)
        if ollama_response.status_code == 200:
            st.success("✅ Ollama Connected")
        else:
            st.error("❌ Ollama Error")
    except:
        st.error("❌ Ollama Disconnected")

    # Check Qdrant connection
    try:
        qdrant_client = QdrantClient(host="localhost", port=6333)
        collections = qdrant_client.get_collections()
        st.success("✅ Qdrant Connected")
        st.info(f"Collections: {len(collections.collections)}")
    except:
        st.error("❌ Qdrant Disconnected")

    # Session info
    if st.session_state.trace_id:
        st.header("📊 Session Info")
        session_status = coordinator.get_session_status(st.session_state.trace_id)
        st.json(session_status)

    # MCP Message History
    st.header("📨 MCP Messages")
    if st.button("Show Last 5 Messages"):
        recent_messages = mcp_bus.message_history[-5:]
        for msg in recent_messages:
            st.text(f"{msg.sender} → {msg.receiver}: {msg.type}")

# File upload section
st.header("📁 Document Upload")
uploaded_files = st.file_uploader(
    "Upload documents (PDF, DOCX, TXT, CSV, PPTX, Images)",
    type=["pdf", "docx", "txt", "csv", "pptx", "png", "jpg", "jpeg"],
    accept_multiple_files=True
)

if uploaded_files:
    # Process uploaded files
    files_data = []

    for file in uploaded_files:
        # Skip if already processed
        if any(f["name"] == file.name for f in st.session_state.files):
            continue

        # Save file to temp directory
        file_path = os.path.join(tempfile.gettempdir(), file.name)
        with open(file_path, "wb") as f:
            f.write(file.read())

        file_info = {
            "name": file.name,
            "path": file_path,
            "ext": file.name.split(".")[-1].lower(),
            "size": os.path.getsize(file_path),
            "uploaded_at": datetime.now().isoformat()
        }

        files_data.append(file_info)
        st.session_state.files.append(file_info)

    if files_data:
        # Process documents using CoordinatorAgent
        with st.spinner("Processing documents..."):
            trace_id = coordinator.process_documents(files_data)
            st.session_state.trace_id = trace_id

            # Wait for processing to complete
            max_wait = 30  # 30 seconds
            wait_time = 0

            while wait_time < max_wait:
                status = coordinator.get_session_status(trace_id)
                if status.get("status") == "ready":
                    st.success(f"✅ Processed {status.get('document_count', 0)} document chunks")
                    break
                elif status.get("status") == "error":
                    st.error(f"❌ Processing error: {status.get('error', 'Unknown error')}")
                    break

                time.sleep(1)
                wait_time += 1

            if wait_time >= max_wait:
                st.warning("⏰ Processing is taking longer than expected...")

# Display uploaded files
if st.session_state.files:
    st.header("📚 Uploaded Documents")

    for i, file_info in enumerate(st.session_state.files):
        with st.expander(f"📄 {file_info['name']} ({file_info['ext'].upper()})"):
            col1, col2 = st.columns(2)

            with col1:
                st.write(f"**Size:** {file_info['size']:,} bytes")
                st.write(f"**Uploaded:** {file_info['uploaded_at']}")

            with col2:
                if st.button(f"🗑️ Remove {file_info['name']}", key=f"remove_{i}"):
                    st.session_state.files.pop(i)
                    st.rerun()

# Query section
if st.session_state.trace_id:
    st.header("💬 Ask Questions")

    # Chat interface
    with st.form("query_form"):
        user_query = st.text_input("🔍 Enter your question:")
        submitted = st.form_submit_button("Ask Question")

        if submitted and user_query:
            with st.spinner("Generating response..."):
                # Process query using CoordinatorAgent
                result = coordinator.process_query(user_query, st.session_state.trace_id)

                # Wait for response
                max_wait = 60
                wait_time = 0

                while wait_time < max_wait:
                    status = coordinator.get_session_status(st.session_state.trace_id)
                    last_response = status.get("last_response")

                    if last_response:
                        # Display response
                        st.subheader("🤖 Response")
                        st.write(last_response["response"])

                        # Show context chunks
                        with st.expander("📖 Source Context"):
                            for j, chunk in enumerate(last_response["context_chunks"]):
                                st.text_area(f"Context {j + 1}", chunk, height=100)

                        # Store in session state
                        st.session_state.query_results[user_query] = last_response
                        break

                    time.sleep(1)
                    wait_time += 1

                if wait_time >= max_wait:
                    st.warning("⏰ Response generation is taking longer than expected...")

# Display previous queries
if st.session_state.query_results:
    st.header("📋 Previous Queries")

    for query, result in st.session_state.query_results.items():
        with st.expander(f"❓ {query}"):
            st.write("**Response:**")
            st.write(result["response"])

            if st.button(f"Clear", key=f"clear_{hash(query)}"):
                del st.session_state.query_results[query]
                st.rerun()

# Footer
st.markdown("---")
st.markdown("*Powered by Agentic Architecture with Model Context Protocol (MCP)*")